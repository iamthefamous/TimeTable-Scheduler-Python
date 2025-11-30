#!/usr/bin/env python3
"""
Generate a PowerPoint presentation for the Timetable Scheduler project.
This script creates a presentation focusing on discrete mathematics concepts.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import os

# Define paths
ASSETS_DIR = os.path.join(os.path.dirname(__file__), 'assets', 'img')
OUTPUT_FILE = os.path.join(os.path.dirname(__file__), 'Timetable_Scheduler_Discrete_Math.pptx')

# Colors
DARK_BLUE = RGBColor(0, 51, 102)
LIGHT_BLUE = RGBColor(0, 102, 204)
ACCENT_GREEN = RGBColor(0, 128, 0)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
GRAY = RGBColor(128, 128, 128)


def add_title_slide(prs):
    """Add the title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add background shape
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "📅 Timetable Scheduler"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Discrete Mathematics & Genetic Algorithm Optimization"
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Add bottom description
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1))
    tf = desc_box.text_frame
    p = tf.paragraphs[0]
    p.text = "An intelligent university timetable generator using evolutionary computation"
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


def add_section_slide(prs, title, subtitle=""):
    """Add a section divider slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add background shape
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = LIGHT_BLUE
    background.line.fill.background()
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER


def add_content_slide(prs, title, bullets, image_path=None):
    """Add a content slide with optional image."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    # Add line under title
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.1),
        Inches(9), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_BLUE
    line.line.fill.background()
    
    # Determine content width based on image presence
    if image_path and os.path.exists(image_path):
        content_width = 5.5
        # Add image on the right
        try:
            slide.shapes.add_picture(image_path, Inches(6.2), Inches(1.5), width=Inches(3.5))
        except Exception as e:
            print(f"Warning: Could not add image {image_path}: {e}")
    else:
        content_width = 9
    
    # Add bullet points
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(content_width), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(18)
        p.font.color.rgb = BLACK
        p.space_before = Pt(8)
        p.space_after = Pt(4)


def add_math_slide(prs, title, content_blocks):
    """Add a slide with mathematical formulas and explanations."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    # Add line under title
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.1),
        Inches(9), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_BLUE
    line.line.fill.background()
    
    # Add content blocks
    y_pos = 1.4
    for block in content_blocks:
        if block.get('type') == 'text':
            box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(0.5))
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = block['content']
            p.font.size = Pt(16)
            p.font.bold = block.get('bold', False)
            p.font.color.rgb = BLACK
            y_pos += 0.4
        elif block.get('type') == 'formula':
            # Formula box with light background
            formula_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.7), Inches(y_pos),
                Inches(8.5), Inches(block.get('height', 0.6))
            )
            formula_shape.fill.solid()
            formula_shape.fill.fore_color.rgb = RGBColor(240, 248, 255)
            formula_shape.line.color.rgb = LIGHT_BLUE
            
            # Add formula text
            formula_box = slide.shapes.add_textbox(
                Inches(0.9), Inches(y_pos + 0.1),
                Inches(8.1), Inches(block.get('height', 0.6) - 0.1)
            )
            tf = formula_box.text_frame
            tf.word_wrap = True
            for i, line in enumerate(block['content'].split('\n')):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(14)
                p.font.name = "Courier New"
                p.font.color.rgb = DARK_BLUE
            y_pos += block.get('height', 0.6) + 0.2


def add_table_slide(prs, title, headers, rows):
    """Add a slide with a table."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    # Add line under title
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.1),
        Inches(9), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_BLUE
    line.line.fill.background()
    
    # Create table
    num_rows = len(rows) + 1  # +1 for header
    num_cols = len(headers)
    
    table_width = Inches(9)
    col_width = table_width / num_cols
    
    table = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.5), Inches(1.4),
        table_width, Inches(0.5 * num_rows)
    ).table
    
    # Style header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    # Add data rows
    for row_idx, row in enumerate(rows):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_text
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = BLACK
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 248, 255)


def add_image_slide(prs, title, image_path, description=""):
    """Add a slide with a large image."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    # Add image if it exists
    if os.path.exists(image_path):
        try:
            slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.2), width=Inches(9))
        except Exception as e:
            print(f"Warning: Could not add image {image_path}: {e}")
    
    # Add description
    if description:
        desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = description
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER


def main():
    """Generate the presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide
    add_title_slide(prs)
    
    # Slide 2: Section - Theoretical Foundations
    add_section_slide(prs, "🔬 Theoretical Foundations", "Discrete Mathematics in Scheduling")
    
    # Slide 3: The Scheduling Problem
    add_content_slide(prs, "The Scheduling Problem as CSP", [
        "Timetable scheduling is a Constraint Satisfaction Problem (CSP)",
        "Variables: A finite set of classes that need to be scheduled",
        "Domains: Discrete set of (room, time slot, instructor) combinations",
        "Constraints: Rules that restrict valid assignments",
        "This is a combinatorial optimization problem",
        "We seek optimal solution from exponentially large set of possibilities"
    ])
    
    # Slide 4: Graph Theory Connections
    add_table_slide(prs, "📊 Graph Theory Connections", 
        ["Graph Theory Concept", "Scheduling Equivalent"],
        [
            ["Vertices", "Classes/Courses to be scheduled"],
            ["Edges", "Conflicts between classes"],
            ["Graph Coloring", "Assigning time slots (colors)"],
            ["Chromatic Number", "Minimum time slots needed"],
            ["NP-hard Problem", "Requires heuristic algorithms"]
        ]
    )
    
    # Slide 5: Computational Complexity
    add_table_slide(prs, "⚡ Computational Complexity",
        ["Aspect", "Complexity Analysis"],
        [
            ["Problem Class", "NP-hard (no known polynomial-time solution)"],
            ["Solution Space", "O((R × T × I)^n) - R=rooms, T=time slots, I=instructors"],
            ["Brute Force", "Infeasible for realistic problem sizes"],
            ["Genetic Algorithm", "Polynomial time per generation, near-optimal solutions"]
        ]
    )
    
    # Slide 6: Set Theory - Mathematical Foundation
    add_math_slide(prs, "📐 Set Theory & Relations",
        [
            {'type': 'text', 'content': 'Formal Definition of the Scheduling Problem:', 'bold': True},
            {'type': 'formula', 'height': 1.8, 'content': 
                "Let S = {s₁, s₂, ..., sₙ} be the set of sections\n"
                "Let T = {t₁, t₂, ..., tₘ} be the set of time slots\n"
                "Let R = {r₁, r₂, ..., rₖ} be the set of rooms\n"
                "Let I = {i₁, i₂, ..., iₗ} be the set of instructors\n"
                "Let C = {c₁, c₂, ..., cₚ} be the set of courses"},
            {'type': 'text', 'content': 'A valid schedule is a function:', 'bold': True},
            {'type': 'formula', 'height': 0.6, 'content': 
                "f: Classes → T × R × I"},
            {'type': 'text', 'content': 'Subject to all constraints being satisfied'},
        ]
    )
    
    # Slide 7: Section - Genetic Algorithm
    add_section_slide(prs, "🧬 Genetic Algorithm", "Discrete Optimization in Action")
    
    # Slide 8: Genetic Algorithm Overview
    add_content_slide(prs, "Genetic Algorithm Overview", [
        "Metaheuristic optimization inspired by natural selection",
        "Operates on principle of 'survival of the fittest'",
        "Effective for discrete optimization problems",
        "Key Operations:",
        "  1. Initialize Population - Random schedules",
        "  2. Fitness Evaluation - Count constraint violations",
        "  3. Selection - Tournament selection",
        "  4. Crossover - Combine parent schedules",
        "  5. Mutation - Random alterations",
        "  6. Termination - Return best schedule"
    ])
    
    # Slide 9: Chromosome Encoding
    add_math_slide(prs, "🧬 Chromosome Encoding",
        [
            {'type': 'text', 'content': 'Each schedule is encoded as a sequence of genes:', 'bold': True},
            {'type': 'formula', 'height': 0.8, 'content': 
                "Schedule = [Gene₁, Gene₂, ..., Geneₙ]\n"
                "where Geneᵢ = (courseᵢ, roomᵢ, timeᵢ, instructorᵢ)"},
            {'type': 'text', 'content': ''},
            {'type': 'text', 'content': 'Fitness Function:', 'bold': True},
            {'type': 'formula', 'height': 0.6, 'content': 
                "fitness(schedule) = 1 / (1 + number_of_conflicts)"},
            {'type': 'text', 'content': ''},
            {'type': 'text', 'content': 'Properties:'},
            {'type': 'text', 'content': '• Perfect schedule (0 conflicts) has fitness = 1.0'},
            {'type': 'text', 'content': '• More conflicts → lower fitness → less likely to survive'},
        ]
    )
    
    # Slide 10: Genetic Operations
    add_table_slide(prs, "⚙️ Genetic Operations",
        ["Operation", "Mathematical Description", "Purpose"],
        [
            ["Selection", "Tournament: max(fitness) from k individuals", "Choose fittest parents"],
            ["Crossover", "Uniform with p=0.5: Gene from A or B", "Combine successful traits"],
            ["Mutation", "Probability p=0.05: Replace with random gene", "Introduce diversity"],
            ["Elitism", "Preserve top 2 schedules unchanged", "Prevent loss of best solutions"]
        ]
    )
    
    # Slide 11: Discrete Math in Genetic Operations
    add_content_slide(prs, "Discrete Math in Genetic Operators", [
        "Tournament Selection uses:",
        "  • Combinatorics: C(n,k) ways to select k from n individuals",
        "  • Probability: Selection proportional to relative fitness",
        "",
        "Crossover applies:",
        "  • Boolean Algebra: Random boolean decides gene parent",
        "  • Permutations: Creating new arrangements",
        "",
        "Mutation employs:",
        "  • Bernoulli trials with p=0.05 for each gene",
        "  • Uniform random sampling from valid domain"
    ])
    
    # Slide 12: Section - Constraints
    add_section_slide(prs, "🔒 Constraints", "Logical Predicates in Discrete Mathematics")
    
    # Slide 13: Hard Constraints as Logic
    add_math_slide(prs, "Hard Constraints as Predicate Logic",
        [
            {'type': 'text', 'content': 'Section Conflict:', 'bold': True},
            {'type': 'formula', 'height': 0.6, 'content': 
                "∀ c₁, c₂ ∈ Classes, c₁ ≠ c₂: ¬(section(c₁) = section(c₂) ∧ time(c₁) = time(c₂))"},
            {'type': 'text', 'content': 'Room Conflict:', 'bold': True},
            {'type': 'formula', 'height': 0.6, 'content': 
                "∀ c₁, c₂ ∈ Classes, c₁ ≠ c₂: ¬(room(c₁) = room(c₂) ∧ time(c₁) = time(c₂))"},
            {'type': 'text', 'content': 'Instructor Conflict:', 'bold': True},
            {'type': 'formula', 'height': 0.6, 'content': 
                "∀ c₁, c₂ ∈ Classes, c₁ ≠ c₂: ¬(instructor(c₁) = instructor(c₂) ∧ time(c₁) = time(c₂))"},
            {'type': 'text', 'content': 'Room Capacity:', 'bold': True},
            {'type': 'formula', 'height': 0.6, 'content': 
                "∀ c ∈ Classes: capacity(room(c)) ≥ students(course(c))"},
        ]
    )
    
    # Slide 14: Conflict Detection as SAT
    add_math_slide(prs, "Conflict Detection as Boolean SAT",
        [
            {'type': 'text', 'content': 'The fitness function transforms this into a weighted MAX-SAT problem:', 'bold': True},
            {'type': 'formula', 'height': 0.6, 'content': 
                "minimize: Σ violation(constraint_i)"},
            {'type': 'formula', 'height': 0.6, 'content': 
                "subject to: all hard constraints should have violation = 0"},
            {'type': 'text', 'content': ''},
            {'type': 'text', 'content': 'Key Insight:', 'bold': True},
            {'type': 'text', 'content': '• Each constraint is a Boolean predicate'},
            {'type': 'text', 'content': '• Satisfied constraint = TRUE (violation = 0)'},
            {'type': 'text', 'content': '• Violated constraint = FALSE (violation = 1)'},
            {'type': 'text', 'content': '• Goal: Find assignment where all predicates are TRUE'},
        ]
    )
    
    # Slide 15: Section - Complexity Analysis
    add_section_slide(prs, "📊 Algorithm Complexity", "Performance Analysis")
    
    # Slide 16: Time Complexity
    add_table_slide(prs, "⏱️ Time Complexity Analysis",
        ["Phase", "Complexity", "Description"],
        [
            ["Initialization", "O(P × n)", "Create P random schedules with n classes"],
            ["Fitness Evaluation", "O(n²)", "Compare all pairs for conflicts"],
            ["Selection", "O(P × k)", "Tournament selection, k size, P times"],
            ["Crossover", "O(P × n)", "Create P offspring with n genes"],
            ["Mutation", "O(P × n)", "Potentially mutate each gene"],
            ["Per Generation", "O(P × n²)", "Dominated by fitness evaluation"],
            ["Total Algorithm", "O(G × P × n²)", "G generations"]
        ]
    )
    
    # Slide 17: Space Complexity and Convergence
    add_content_slide(prs, "Space Complexity & Convergence", [
        "Space Complexity:",
        "  • Population: O(P × n) - Store P schedules with n classes",
        "  • Auxiliary: O(n) - Temporary storage during operations",
        "  • Total: O(P × n) - Linear in population and problem size",
        "",
        "Convergence Properties:",
        "  • Exploration vs Exploitation: Mutation rate (5%) balances both",
        "  • Elitism: Ensures monotonic improvement of best fitness",
        "  • Tournament Selection: Creates selection pressure",
        "",
        "Algorithm terminates when:",
        "  • Perfect fitness (1.0) achieved, or",
        "  • Maximum generations (100) reached"
    ])
    
    # Slide 18: Algorithm Parameters
    add_table_slide(prs, "🎛️ Algorithm Parameters",
        ["Parameter", "Value", "Mathematical Rationale"],
        [
            ["Population Size", "30", "Balance between diversity and computation"],
            ["Elite Schedules", "2", "Guarantees monotonic improvement"],
            ["Tournament Size", "8", "~27% of population for selection pressure"],
            ["Mutation Rate", "0.05", "5% chance per gene for exploration"],
            ["Max Generations", "100", "Convergence typically within 50-100"]
        ]
    )
    
    # Slide 19: Section - Application
    add_section_slide(prs, "💻 Application Screenshots", "Django Web Application")
    
    # Add screenshot slides
    images = [
        ("1.png", "Home Page", "Main dashboard for the timetable scheduler"),
        ("2.png", "Add Instructor", "Interface for adding new instructors"),
        ("3.png", "Add Room", "Room management with capacity settings"),
        ("4.png", "Add Course", "Course creation and instructor assignment"),
        ("7.png", "Generated Timetable", "Optimized schedule produced by the genetic algorithm"),
        ("8.png", "Timetable View", "Final timetable display with all assignments")
    ]
    
    for img_file, title, desc in images:
        img_path = os.path.join(ASSETS_DIR, img_file)
        if os.path.exists(img_path):
            add_image_slide(prs, title, img_path, desc)
    
    # Slide: Related Problems
    add_section_slide(prs, "📚 Related Problems", "Classical Combinatorial Optimization")
    
    add_content_slide(prs, "Related Discrete Math Problems", [
        "Graph Coloring Problem:",
        "  • Assign colors to vertices with no adjacent same colors",
        "  • Scheduling = coloring with time slots as colors",
        "",
        "Bin Packing Problem:",
        "  • Fit items into minimum containers",
        "  • Scheduling = packing classes into time slots",
        "",
        "Job Shop Scheduling:",
        "  • Schedule jobs on machines with precedence",
        "  • Similar constraint satisfaction structure",
        "",
        "All these are NP-hard → Require heuristic approaches"
    ])
    
    # Final Slide: Summary
    add_section_slide(prs, "🎯 Key Takeaways", "")
    
    add_content_slide(prs, "Summary: Discrete Math in Scheduling", [
        "✅ CSP Framework: Variables, Domains, Constraints",
        "✅ Graph Theory: Scheduling ≡ Graph Coloring (NP-hard)",
        "✅ Set Theory: Formal definitions using sets and relations",
        "✅ Predicate Logic: Constraints as Boolean predicates",
        "✅ Combinatorics: Exponential solution space O((R×T×I)^n)",
        "✅ Genetic Algorithm: Polynomial-time heuristic approach",
        "✅ Optimization: MAX-SAT formulation for constraint satisfaction",
        "",
        "The power of discrete mathematics transforms complex",
        "real-world scheduling into solvable algorithmic problems!"
    ])
    
    # Thank You Slide
    add_section_slide(prs, "Thank You! 🙏", "Questions?")
    
    # Save the presentation
    prs.save(OUTPUT_FILE)
    print(f"Presentation saved to: {OUTPUT_FILE}")


if __name__ == "__main__":
    main()
